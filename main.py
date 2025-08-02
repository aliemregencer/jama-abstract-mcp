#!/usr/bin/env python3
"""
JAMA Abstract MCP Server
JAMA tıp dergisi makalelerinden veri çıkarma MCP servisi
"""

import asyncio
import json
import logging
import os
import tempfile
from typing import Any, Dict, List, Optional
from io import BytesIO

# FastMCP import
import fastmcp

# PPTX ve görsel işleme kütüphaneleri
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch
import seaborn as sns
import numpy as np
from PIL import Image, ImageDraw, ImageFont

from scraper import JAMAScraper
from parser import DataParser

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# FastMCP server instance
mcp = fastmcp.FastMCP("jama-abstract-mcp")

@mcp.tool()
async def extract_jama_article(url: str) -> Dict[str, Any]:
    """
    JAMA makale linkinden makale verilerini çıkarır
    
    Args:
        url: JAMA makale URL'si
        
    Returns:
        Çıkarılan makale verileri (başlık, yazarlar, özet, vb.)
    """
    return await _extract_jama_article(url)

async def _extract_jama_article(url: str) -> Dict[str, Any]:
    try:
        logger.info(f"JAMA makale çıkarma başlıyor: {url}")
        
        # URL doğrulama
        if "jamanetwork.com" not in url:
            return {"error": "Geçersiz JAMA URL'si"}
        
        # Scraper ile veri çıkarma
        scraper = JAMAScraper(headless=True, timeout=8)
        html_content = await scraper.scrape_article(url)
        
        if not html_content:
            return {"error": "Makale içeriği alınamadı"}
        
        # Parser ile veri ayrıştırma
        parser = DataParser()
        article_data = parser.parse_article(html_content)
        
        logger.info("Makale verisi başarıyla çıkarıldı")
        return {
            "success": True,
            "data": article_data,
            "source_url": url
        }
        
    except Exception as e:
        logger.error(f"extract_jama_article hatası: {e}")
        return {"error": f"Veri çıkarma hatası: {str(e)}"}

@mcp.tool()
async def create_abstract_visual(article_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    Makale verilerine göre abstract görsel oluşturur ve PPTX formatında döndürür
    
    Args:
        article_data: extract_jama_article tool'undan gelen makale verileri
        
    Returns:
        PPTX dosyası (base64 encoded) ve metadata
    """
    return await _create_abstract_visual(article_data)

async def _create_abstract_visual(article_data: Dict[str, Any]) -> Dict[str, Any]:
    try:
        logger.info("Abstract görsel oluşturma başlıyor")
        
        if not article_data or "data" not in article_data:
            return {"error": "Geçersiz makale verisi"}
        
        data = article_data["data"]
        
        # PPTX oluştur
        pptx_bytes = await _generate_abstract_pptx(data)
        
        if pptx_bytes:
            logger.info("Abstract görsel başarıyla oluşturuldu")
            return {
                "success": True,
                "pptx_base64": pptx_bytes,
                "filename": f"jama_abstract_{data.get('title', 'study').replace(' ', '_')[:50]}.pptx",
                "message": "Abstract görsel PPTX formatında hazırlandı"
            }
        else:
            return {"error": "PPTX oluşturulamadı"}
        
    except Exception as e:
        logger.error(f"create_abstract_visual hatası: {e}")
        return {"error": f"Görsel oluşturma hatası: {str(e)}"}

async def _generate_abstract_pptx(article_data: Dict[str, Any]) -> str:
    """Makale verilerine göre abstract görsel PPTX oluşturur"""
    try:
        # Yeni presentation oluştur
        prs = Presentation()
        
        # Slide boyutunu ayarla (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Slide ekle
        slide_layout = prs.slide_layouts[6]  # Boş slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Arka plan rengini ayarla
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Başlık ekle
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = article_data.get("title", "Study Title")
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_font = title_para.font
        title_font.size = Pt(24)
        title_font.bold = True
        title_font.color.rgb = RGBColor(0, 0, 0)
        
        # Ana içerik alanları
        content_width = Inches(3.5)
        content_height = Inches(4)
        
        # Sol üst - Population
        population_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), content_width, content_height)
        population_frame = population_box.text_frame
        population_frame.text = f"POPULATION\n\n{article_data.get('population', 'N/A')}"
        population_para = population_frame.paragraphs[0]
        population_para.font.bold = True
        population_para.font.size = Pt(16)
        population_para.font.color.rgb = RGBColor(0, 0, 139)
        
        # Orta üst - Intervention
        intervention_box = slide.shapes.add_textbox(Inches(4.5), Inches(1.8), content_width, content_height)
        intervention_frame = intervention_box.text_frame
        intervention_frame.text = f"INTERVENTION\n\n{article_data.get('intervention', 'N/A')}"
        intervention_para = intervention_frame.paragraphs[0]
        intervention_para.font.bold = True
        intervention_para.font.size = Pt(16)
        intervention_para.font.color.rgb = RGBColor(0, 0, 139)
        
        # Sağ üst - Findings
        findings_box = slide.shapes.add_textbox(Inches(8.5), Inches(1.8), content_width, content_height)
        findings_frame = findings_box.text_frame
        findings_frame.text = f"FINDINGS\n\n{article_data.get('findings', 'N/A')}"
        findings_para = findings_frame.paragraphs[0]
        findings_para.font.bold = True
        findings_para.font.size = Pt(16)
        findings_para.font.color.rgb = RGBColor(0, 0, 139)
        
        # Sol alt - Settings
        settings_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), content_width, Inches(1))
        settings_frame = settings_box.text_frame
        settings_frame.text = f"SETTINGS\n\n{article_data.get('settings', 'N/A')}"
        settings_para = settings_frame.paragraphs[0]
        settings_para.font.bold = True
        settings_para.font.size = Pt(14)
        settings_para.font.color.rgb = RGBColor(0, 0, 139)
        
        # Orta alt - Primary Outcome
        outcome_box = slide.shapes.add_textbox(Inches(4.5), Inches(6.2), content_width, Inches(1))
        outcome_frame = outcome_box.text_frame
        outcome_frame.text = f"PRIMARY OUTCOME\n\n{article_data.get('primary_outcome', 'N/A')}"
        outcome_para = outcome_frame.paragraphs[0]
        outcome_para.font.bold = True
        outcome_para.font.size = Pt(14)
        outcome_para.font.color.rgb = RGBColor(0, 0, 139)
        
        # PPTX'i bytes'a çevir
        pptx_stream = BytesIO()
        prs.save(pptx_stream)
        pptx_bytes = pptx_stream.getvalue()
        
        # Base64 encode
        import base64
        return base64.b64encode(pptx_bytes).decode('utf-8')
        
    except Exception as e:
        logger.error(f"PPTX oluşturma hatası: {e}")
        return None

if __name__ == "__main__":
    import argparse
    import os

    parser = argparse.ArgumentParser(description="JAMA Abstract MCP Server")
    parser.add_argument("--transport", choices=["stdio", "http"], default="stdio",
                       help="Transport protocol to use")
    parser.add_argument("--host", default="127.0.0.1",
                       help="Host to bind to (for HTTP transport)")
    parser.add_argument("--port", type=int, default=8000,
                       help="Port to bind to (for HTTP transport)")

    args = parser.parse_args()

    # Environment variable kontrolü
    transport = os.getenv("MCP_TRANSPORT", args.transport)
    host = os.getenv("MCP_HOST", args.host)
    port = int(os.getenv("MCP_PORT", args.port))

    if transport == "http":
        logger.info(f"Starting HTTP server on {host}:{port}")
        mcp.run(transport="http", host=host, port=port)
    else:
        logger.info("Starting STDIO server")
        mcp.run()